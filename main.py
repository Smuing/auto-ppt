# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
import copy
import requests
import olefile
import zlib
import struct

def file_reader(filename):
    f = olefile.OleFileIO(filename)
    dirs = f.listdir()

    # 문서 포맷 압축 여부 확인
    header = f.openstream("FileHeader")
    header_data = header.read()
    is_compressed = (header_data[36] & 1) == 1

    # Body Sections 불러오기
    nums = []
    for d in dirs:
        if d[0] == "BodyText":
            nums.append(int(d[1][len("Section"):]))
    sections = ["BodyText/Section" + str(x) for x in sorted(nums)]

    # 전체 text 추출
    text = []
    for section in sections:
        bodytext = f.openstream(section)
        data = bodytext.read()
        if is_compressed:
            unpacked_data = zlib.decompress(data, -15)
        else:
            unpacked_data = data

        # 각 Section 내 text 추출
        section_text = []
        i = 0
        size = len(unpacked_data)
        while i < size:
            header = struct.unpack_from("<I", unpacked_data, i)[0]
            rec_type = header & 0x3ff
            rec_len = (header >> 20) & 0xfff

            if rec_type in [67]:
                rec_data = unpacked_data[i + 4:i + 4 + rec_len - 4]
                try:
                    verseFilter = int(rec_data.decode('utf-16').split(maxsplit=2)[1].split(':')[0])
                    sp = rec_data.decode('utf-16').split(maxsplit=2)
                    section_text.append(sp[0] + ' ' + sp[1] + '\n' + sp[2])
                except:
                    section_text = section_text

            i += 4 + rec_len

        text = section_text

    print_ppt(text, '3')

def api(userInput):
    bible = []
    for i in userInput:
        doc = { "창세기" : "ge", "출애굽기" :  "exo", "레위기" :  "lev", "민수기" :  "num", "신명기" :  "deu", "여호수아" :  "josh", "사사기" :  "jdgs", "룻기" :  "ruth", "사무엘상" : "1sm", "사무엘하" : "2sm", "열왕기상" : "1ki", "열왕기하" : "2ki", "역대상" : "1chr", "역대하" : "2chr", "에스라" :  "ezra", "느헤미야" :  "neh", "에스더" :  "est", "욥기" :  "job", "시편" :  "psa", "잠언" :  "prv", "전도서" :  "eccl", "아가" :  "ssol", "이사야" :  "isa", "예레미야" :  "jer", "예레미야애가" : "lam", "에스겔" :  "eze", "다니엘" :  "dan", "호세아" :  "hos", "요엘" :  "joel", "아모스" :  "amos", "오바댜" :  "obad", "요나" :  "jonah", "미가" :  "mic", "나훔" :  "nahum", "하박국" :  "hab", "스바냐" :  "zep", "학개" :  "hag", "스가랴" :  "zep", "말라기" :  "mal", "마태복음" :  "mat", "마가복음" : "mark", "누가복음" : "luke", "요한복음" : "john", "사도행전" : "acts", "로마서" : "rom", "고린도전서" : "1cor", "고린도후서" : "2cor", "갈라디아서" : "gal", "에베소서" : "eph", "빌립보서" : "phi", "골로새서" : "col", "데살로니가전서" : "1th", "데살로니가후서" : "2th", "디모데전서" : "1tim", "디모데후서" : "2tim", "디도서" : "titus", "빌레몬서" : "phmn", "히브리서" : "heb", "야고보서" : "jas", "베드로전서" : "1pet", "베드로후서" : "2pet", "요한1서" : "1jn", "요한2서" : "2jn", "요한3서" : "3jn", "유다서" : "jude", "요한계시록" : "rev",
                "창" : "ge", "출" :  "exo", "레" :  "lev", "민" :  "num", "신" :  "deu", "수" :  "josh", "삿" :  "jdgs", "룻" :  "ruth", "삼상" : "1sm", "삼하" : "2sm", "왕상" : "1ki", "왕하" : "2ki", "대상" : "1chr", "대하" : "2chr", "스" :  "ezra", "느" :  "neh", "에" :  "est", "욥" :  "job", "시" :  "psa", "잠" :  "prv", "전" :  "eccl", "아" :  "ssol", "사" :  "isa", "렘" :  "jer", "애" : "lam", "겔" :  "eze", "단" :  "dan", "호" :  "hos", "욜" :  "joel", "암" :  "amos", "옵" :  "obad", "욘" :  "jonah", "미" :  "mic", "나" :  "nahum", "합" :  "hab", "습" :  "zep", "학" :  "hag", "슥" :  "zep", "말" :  "mal", "마" :  "mat", "막" : "mark", "눅" : "luke", "요" : "john", "행" : "acts", "롬" : "rom", "고전" : "1cor", "고후" : "2cor", "갈" : "gal", "엡" : "eph", "빌" : "phi", "골" : "col", "살전" : "1th", "살후" : "2th", "딤전" : "1tim", "딤후" : "2tim", "딛" : "titus", "몬" : "phmn", "히" : "heb", "약" : "jas", "벧전" : "1pet", "벧후" : "2pet", "요일" : "1jn", "요이" : "2jn", "요삼" : "3jn", "유" : "jude", "계" : "rev"}\
                .get(i['doc'])
        if i['type'] == '1':
            res = requests.get(
                'https://yesu.io/bible?lang=kor&doc=' + doc + '&start=' + i['start'] + '&end=' + i['end'])
            result = eval(res.text)
            result[0].update({"doc": i['doc']})
            bible = result
            print_ppt(bible, '1')

        elif i['type'] == '2':
            res = requests.get('https://yesu.io/bible?lang=kor&doc='+doc+'&start='+i['verse']+'&end='+i['verse'])
            result = eval(res.text)
            result[0].update({"doc": i['doc']})
            bible.append(result[0])
            print_ppt(bible, '2')

def print_ppt(bible, userType):
    pptx_fpath = './theme.pptx'
    prs = Presentation(pptx_fpath)

    if userType == '3':
        for i, b in enumerate(bible):
            source_slide = prs.slides[0]
            if i != 0:
                slide_layout = prs.slide_layouts[6]
                copied_slide = prs.slides.add_slide(slide_layout)
                el = source_slide.shapes[0].element
                newel = copy.deepcopy(el)
                copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
                shape = copied_slide.shapes[0]
            else:
                shape = source_slide.shapes[0]
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = b
            font = run.font
            font.color.rgb = RGBColor(0, 0, 0)
            font.name = 'HY견고딕'
            font.size = Pt(33)
    else:
        for i, b in enumerate(bible):
            source_slide = prs.slides[0]
            if i != 0:
                slide_layout = prs.slide_layouts[6]
                copied_slide = prs.slides.add_slide(slide_layout)
                el = source_slide.shapes[0].element
                newel = copy.deepcopy(el)
                copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
                shape = copied_slide.shapes[0]
                if userType == '1':
                    t = b['verse'] + '. ' + b['message']
                elif userType == '2':
                    t = b['doc'] + ' ' + b['chapter'] + ':' + b['verse'] + '\n' + b['message']
            else:
                shape = source_slide.shapes[0]
                if userType == '1':
                    t = b['doc'] + ' ' + b['chapter'] + '장\n' + b['verse'] + '. ' + b['message']
                elif userType == '2':
                    t = b['doc'] + ' ' + b['chapter'] + ':' + b['verse'] + '\n' + b['message']
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = t
            font = run.font
            font.color.rgb = RGBColor(0, 0, 0)
            font.name = 'HY견고딕'
            font.size = Pt(33)

    prs.save('result.pptx')

if __name__ == '__main__':
    userInput = []
    des = "봉독 > 1\n설교 말씀 > 2\n설교 한글파일 > 3\n그만 입력 > e\n> "
    while True:
        inputType = input(des)
        if inputType == '1':
            if des == "설교 말씀 > 2\n그만 입력하기 > e\n> ":
                continue
            else:
                doc = input("문서 : ")
                start = input("시작 장:절 : ")
                end = input("끝 장:절 : ")
                userInput.append({"type": inputType, "doc": doc, "start": start, "end": end})
                api(userInput)
                break
        elif inputType == '2':
            doc = input("문서 : ")
            verse = input("장:절 : ")
            userInput.append({"type": inputType, "doc": doc, "verse": verse})
            des = "설교 말씀 > 2\n그만 입력하기 > e\n> "
            api(userInput)
        elif inputType == '3':
            fileName = input("파일명 : ")
            file_reader(fileName + '.hwp')
            break
        elif inputType == 'e':
            break
        else:
            continue